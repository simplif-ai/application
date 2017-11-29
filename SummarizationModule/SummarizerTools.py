"""
AUTHORS: Kevin Xia

PURPOSE:
    Used for summarizing strange file formats.

DEVELOPER NOTES:
    None
"""

# =============================================================================
# Libraries and Global Variables
# =============================================================================

import sys
sys.path.append('../')

from pptx import Presentation

# =============================================================================


class SummarizerTools:
    """ Summarizer Tool objects generated to deal with strange file formats """

    def __init__(self):
        """ Initialize field variables """

    def extract_ppt(self, filename):
        """ Extracts text from ppt files """
        prs = Presentation(filename)

        sents = []
        for slide in prs.slides:
            for shape in slide.shapes:
                sents.append(shape.text)

        text = ""
        for sent in sents:
            for bullet in sent.split('\n'):
                bullstr = bullet.strip()
                if len(bullstr) > 0:
                    text += bullstr
                    if bullstr[-1] != '.' and bullstr[-1] != '!' and bullstr[-1] != '?':
                        text += '.'
                    text += ' '

        return text
